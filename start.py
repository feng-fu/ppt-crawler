from pptx import Presentation
from pptx.util import Inches
import codecs

images = ['jpg', 'png', 'gif']

def getPaths(path):
    movies_info = []

    with codecs.open(path, 'r', 'utf-8') as f:
        for line in f.readlines():
          data = {}
          infos = line.replace('\n', '').split('|')
          data['url'] = infos[0]
          data['poster'] = infos[1]
          movies_info.append(data)
    return movies_info


def generPpt(paths):

  prs = Presentation()

  for i, path in enumerate(paths):
    url = path['url']
    poster = path['poster']
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = top = Inches(0)
    width = Inches(10)
    height = Inches(7.5)
    fileType = url.split('.')[-1]
    if fileType == "jpg":
      slide.shapes.add_picture(url, left, top)
    else:
      slide.shapes.add_movie(url, left, top, width, height, poster_frame_image=poster, mime_type='video/' + fileType)
  prs.save("./ppt/Lesson2-What a day!.pptx")

# left = Inches(5)
# height = Inches(5.5)
# pic = slide.shapes.add_picture(img_path, left, top, height=height)
def main():
    paths = getPaths('./pptImages.txt')
    generPpt(paths)

if __name__ == '__main__':
    main()
